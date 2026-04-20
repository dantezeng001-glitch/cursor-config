from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import unicodedata
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Any


DEFAULT_CONFIG = {
    "source_patterns": ["*.pdf", "*.ppt", "*.pptx", "**/*.pdf", "**/*.ppt", "**/*.pptx"],
    "exclude_patterns": ["cursor_docs/**", ".cursor/**", ".vscode/**", "**/~$*.pptx", "~$*.pptx"],
    "output_dir": "cursor_docs",
    "pdf_preview_pages": 8,
    "ppt_preview_slides": 12,
    "enable_ocr_fallback": True,
    "ocr_languages": "chi_sim+eng",
    "convert_legacy_ppt": True,
}

CONFIG_FILE_NAME = "document_sync_config.json"


def load_config(project_root: Path) -> dict[str, Any]:
    config_path = project_root / CONFIG_FILE_NAME
    config = DEFAULT_CONFIG.copy()
    if config_path.exists():
        user_config = json.loads(config_path.read_text(encoding="utf-8"))
        config.update(user_config)
    return config


def is_excluded(relative_path: str, patterns: list[str]) -> bool:
    path = PurePosixPath(relative_path)
    return any(path.match(pattern) for pattern in patterns)


def discover_documents(project_root: Path, config: dict[str, Any]) -> list[Path]:
    candidates: set[Path] = set()
    for pattern in config["source_patterns"]:
        candidates.update(project_root.glob(pattern))

    documents: list[Path] = []
    for path in sorted(candidates):
        if not path.is_file():
            continue
        if path.suffix.lower() not in {".pdf", ".ppt", ".pptx"}:
            continue
        relative_path = path.relative_to(project_root).as_posix()
        if is_excluded(relative_path, config["exclude_patterns"]):
            continue
        if path.name.startswith("~$"):
            continue
        documents.append(path)
    return documents


def safe_name(name: str) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*]+', "_", name).strip()
    cleaned = re.sub(r"\s+", "_", cleaned)
    return cleaned or "untitled"


def document_id(project_root: Path, document_path: Path) -> str:
    relative_path = document_path.relative_to(project_root).as_posix()
    suffix = hashlib.md5(relative_path.encode("utf-8")).hexdigest()[:8]
    return f"{safe_name(document_path.stem)}_{suffix}"


def normalize_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in text.split("\n")]
    return "\n".join(line for line in lines if line)


def write_text(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def write_json(path: Path, payload: Any) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def convert_legacy_ppt(document_path: Path, working_dir: Path) -> Path | None:
    helper_path = Path(__file__).with_name("convert_legacy_ppt.ps1")
    converted_path = working_dir / f"{safe_name(document_path.stem)}.converted.pptx"
    result = subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(helper_path),
            "-InputPath",
            str(document_path),
            "-OutputPath",
            str(converted_path),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not converted_path.exists():
        return None
    return converted_path


def ocr_pdf_page(pdf_path: Path, page_index: int, languages: str) -> str:
    try:
        import fitz
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    document = fitz.open(str(pdf_path))
    try:
        page = document.load_page(page_index)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        return normalize_text(pytesseract.image_to_string(image, lang=languages))
    except Exception:
        return ""
    finally:
        document.close()


def _extract_pdf_pymupdf(document_path: Path, config: dict[str, Any]) -> tuple[dict[str, Any], str] | None:
    """Try extracting PDF text with pymupdf (better CJK support). Returns None if unavailable."""
    try:
        import fitz  # pymupdf
    except ImportError:
        return None

    doc = fitz.open(str(document_path))
    pages: list[dict[str, Any]] = []
    full_text_chunks: list[str] = []
    ocr_pages = 0

    for index in range(len(doc)):
        page = doc.load_page(index)
        text = normalize_text(page.get_text("text") or "")
        ocr_used = False
        if not text and config["enable_ocr_fallback"]:
            text = ocr_pdf_page(document_path, index, config["ocr_languages"])
            ocr_used = bool(text)
            if ocr_used:
                ocr_pages += 1

        pages.append({"page": index + 1, "text": text, "ocr_used": ocr_used})
        full_text_chunks.append(f"## Page {index + 1}\n\n{text or '[No text extracted]'}")

    doc.close()

    payload = {
        "type": "pdf",
        "page_count": len(pages),
        "ocr_pages": ocr_pages,
        "pages": pages,
    }
    return payload, "\n\n".join(full_text_chunks).strip() + "\n"


def extract_pdf(document_path: Path, config: dict[str, Any]) -> tuple[dict[str, Any], str]:
    result = _extract_pdf_pymupdf(document_path, config)
    if result is not None:
        return result

    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Missing required dependency: pypdf (or pymupdf)") from exc

    reader = PdfReader(str(document_path))
    pages: list[dict[str, Any]] = []
    full_text_chunks: list[str] = []
    ocr_pages = 0

    for index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        ocr_used = False
        if not text and config["enable_ocr_fallback"]:
            text = ocr_pdf_page(document_path, index - 1, config["ocr_languages"])
            ocr_used = bool(text)
            if ocr_used:
                ocr_pages += 1

        pages.append({"page": index, "text": text, "ocr_used": ocr_used})
        full_text_chunks.append(f"## Page {index}\n\n{text or '[No text extracted]'}")

    payload = {
        "type": "pdf",
        "page_count": len(pages),
        "ocr_pages": ocr_pages,
        "pages": pages,
    }
    return payload, "\n\n".join(full_text_chunks).strip() + "\n"


def shape_lines(shape: Any) -> list[str]:
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = normalize_text(shape.text or "")
        if text:
            lines.extend(text.split("\n"))
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_text = " | ".join(normalize_text(cell.text) for cell in row.cells)
            if row_text.strip(" |"):
                lines.append(row_text)
    return lines


def extract_pptx(document_path: Path) -> tuple[dict[str, Any], str]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing required dependency: python-pptx") from exc

    presentation = Presentation(str(document_path))
    slides_payload: list[dict[str, Any]] = []
    markdown_lines: list[str] = [f"# {document_path.name}", ""]

    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = normalize_text(slide.shapes.title.text or "")

        content_lines: list[str] = []
        for shape in slide.shapes:
            content_lines.extend(shape_lines(shape))

        content_lines = [line for line in content_lines if line and line != title]
        slide_title = title or (content_lines[0][:80] if content_lines else f"Slide {index}")
        slides_payload.append({"slide": index, "title": slide_title, "content": content_lines})

        markdown_lines.append(f"## Slide {index}: {slide_title}")
        markdown_lines.append("")
        if content_lines:
            markdown_lines.extend(f"- {line}" for line in content_lines)
        else:
            markdown_lines.append("- [No text extracted]")
        markdown_lines.append("")

    payload = {"type": "pptx", "slide_count": len(slides_payload), "slides": slides_payload}
    return payload, "\n".join(markdown_lines).rstrip() + "\n"


def write_pdf_markdown(path: Path, document_name: str, payload: dict[str, Any], preview_pages: int) -> None:
    lines = [
        f"# {document_name}",
        "",
        f"- Type: `PDF`",
        f"- Pages: `{payload['page_count']}`",
        f"- OCR pages: `{payload['ocr_pages']}`",
        "",
        "## Preview",
        "",
    ]

    for page in payload["pages"][:preview_pages]:
        snippet = page["text"][:1200] if page["text"] else "[No text extracted]"
        lines.extend(
            [
                f"### Page {page['page']}",
                "",
                snippet,
                "",
            ]
        )

    if payload["page_count"] > preview_pages:
        lines.append(f"Only the first `{preview_pages}` pages are shown here. Use the sibling `.txt` for full extracted text.")

    write_text(path, "\n".join(lines).rstrip() + "\n")


def write_ppt_markdown(path: Path, document_name: str, payload: dict[str, Any], preview_slides: int) -> None:
    lines = [
        f"# {document_name}",
        "",
        f"- Type: `Presentation`",
        f"- Slides: `{payload['slide_count']}`",
        "",
        "## Preview",
        "",
    ]

    for slide in payload["slides"][:preview_slides]:
        lines.append(f"### Slide {slide['slide']}: {slide['title']}")
        lines.append("")
        if slide["content"]:
            lines.extend(f"- {line}" for line in slide["content"][:20])
        else:
            lines.append("- [No text extracted]")
        lines.append("")

    if payload["slide_count"] > preview_slides:
        lines.append(f"Only the first `{preview_slides}` slides are shown here. Use the sibling `.json` for full structured content.")

    write_text(path, "\n".join(lines).rstrip() + "\n")


def sync_document(project_root: Path, output_dir: Path, document_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    doc_id = document_id(project_root, document_path)
    doc_output_dir = output_dir / doc_id
    if doc_output_dir.exists():
        shutil.rmtree(doc_output_dir)
    doc_output_dir.mkdir(parents=True, exist_ok=True)

    relative_source = document_path.relative_to(project_root).as_posix()
    suffix = document_path.suffix.lower()
    manifest: dict[str, Any] = {
        "id": doc_id,
        "name": document_path.name,
        "source": relative_source,
        "updated_at": datetime.fromtimestamp(document_path.stat().st_mtime).isoformat(timespec="seconds"),
        "type": suffix.lstrip("."),
    }

    if suffix == ".pdf":
        payload, full_text = extract_pdf(document_path, config)
        md_path = doc_output_dir / "document.md"
        txt_path = doc_output_dir / "document.txt"
        json_path = doc_output_dir / "document.json"
        write_pdf_markdown(md_path, document_path.name, payload, config["pdf_preview_pages"])
        write_text(txt_path, full_text)
        write_json(json_path, payload)
        manifest["artifacts"] = {
            "md": md_path.relative_to(project_root).as_posix(),
            "txt": txt_path.relative_to(project_root).as_posix(),
            "json": json_path.relative_to(project_root).as_posix(),
        }
        manifest["page_count"] = payload["page_count"]
        manifest["ocr_pages"] = payload["ocr_pages"]
        return manifest

    if suffix == ".ppt":
        converted = None
        if config["convert_legacy_ppt"]:
            converted = convert_legacy_ppt(document_path, doc_output_dir)
        if converted is None:
            md_path = doc_output_dir / "document.md"
            json_path = doc_output_dir / "document.json"
            write_text(
                md_path,
                f"# {document_path.name}\n\nLegacy `.ppt` file detected.\n\nAutomatic conversion failed. Install Microsoft PowerPoint or convert this file to `.pptx` manually.\n",
            )
            write_json(
                json_path,
                {
                    "type": "ppt",
                    "status": "conversion_failed",
                    "message": "PowerPoint automation is unavailable or conversion failed.",
                },
            )
            manifest["artifacts"] = {
                "md": md_path.relative_to(project_root).as_posix(),
                "json": json_path.relative_to(project_root).as_posix(),
            }
            manifest["status"] = "conversion_failed"
            return manifest
        document_path = converted

    payload, markdown_text = extract_pptx(document_path)
    md_path = doc_output_dir / "document.md"
    json_path = doc_output_dir / "document.json"
    write_ppt_markdown(md_path, document_path.name, payload, config["ppt_preview_slides"])
    write_json(json_path, payload)
    write_text(doc_output_dir / "document_full.md", markdown_text)
    manifest["artifacts"] = {
        "md": md_path.relative_to(project_root).as_posix(),
        "full_md": (doc_output_dir / "document_full.md").relative_to(project_root).as_posix(),
        "json": json_path.relative_to(project_root).as_posix(),
    }
    manifest["slide_count"] = payload["slide_count"]
    if suffix == ".ppt":
        manifest["converted_from_ppt"] = True
    return manifest


def cleanup_stale_outputs(output_dir: Path, active_ids: set[str]) -> None:
    if not output_dir.exists():
        return
    for child in output_dir.iterdir():
        if child.is_dir() and child.name not in active_ids:
            shutil.rmtree(child)


def build_index_markdown(index_path: Path, manifest: list[dict[str, Any]]) -> None:
    lines = [
        "# Document Sync Index",
        "",
        "Read this file first when a request is about PDFs, slides, decks, presentations, brochures, or spec documents.",
        "",
        f"- Updated at: `{datetime.now().isoformat(timespec='seconds')}`",
        f"- Documents tracked: `{len(manifest)}`",
        "",
        "## Documents",
        "",
    ]

    if not manifest:
        lines.append("No PDF or presentation files are currently being tracked.")

    for item in manifest:
        lines.extend(
            [
                f"### {item['name']}",
                f"- Source: `{item['source']}`",
                f"- Type: `{item['type']}`",
                f"- Updated at: `{item['updated_at']}`",
            ]
        )
        for key, value in item.get("artifacts", {}).items():
            lines.append(f"- {key}: `{value}`")
        if "page_count" in item:
            lines.append(f"- Pages: `{item['page_count']}`")
        if "slide_count" in item:
            lines.append(f"- Slides: `{item['slide_count']}`")
        if "ocr_pages" in item:
            lines.append(f"- OCR pages: `{item['ocr_pages']}`")
        if "status" in item:
            lines.append(f"- Status: `{item['status']}`")
        lines.append("")

    write_text(index_path, "\n".join(lines).rstrip() + "\n")


def run_sync(project_root: Path) -> list[dict[str, Any]]:
    config = load_config(project_root)
    output_dir = project_root / config["output_dir"]
    output_dir.mkdir(parents=True, exist_ok=True)

    documents = discover_documents(project_root, config)
    manifests = [sync_document(project_root, output_dir, document_path, config) for document_path in documents]
    cleanup_stale_outputs(output_dir, {item["id"] for item in manifests})

    write_json(output_dir / "index.json", manifests)
    build_index_markdown(output_dir / "index.md", manifests)
    return manifests


def _locate_md_cleanup_script() -> Path | None:
    """Best-effort lookup for the md-cleanup skill's cleanup script.

    Search order:
      1. $CURSOR_HOME/skills/md-cleanup/scripts/md_cleanup.py
      2. ~/.cursor/skills/md-cleanup/scripts/md_cleanup.py

    Returns None if the skill isn't installed; callers should silently skip.
    """
    candidates: list[Path] = []
    cursor_home = os.environ.get("CURSOR_HOME")
    if cursor_home:
        candidates.append(Path(cursor_home) / "skills" / "md-cleanup" / "scripts" / "md_cleanup.py")
    candidates.append(Path.home() / ".cursor" / "skills" / "md-cleanup" / "scripts" / "md_cleanup.py")
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    return None


def run_md_cleanup(output_dir: Path) -> None:
    """Invoke the md-cleanup skill's script on every `.md` under output_dir.

    Best-effort: prints a diagnostic line if md-cleanup is missing, but never
    raises — the sync itself already succeeded.
    """
    script = _locate_md_cleanup_script()
    if script is None:
        print("  [md-cleanup] skill not found; skipping post-processing.")
        print("              Install: https://github.com/.../md-cleanup (or see file-sync/SKILL.md).")
        return

    md_files = sorted(str(p) for p in output_dir.rglob("*.md"))
    if not md_files:
        return

    try:
        subprocess.run(
            [sys.executable, str(script), *md_files, "--in-place"],
            check=False,
        )
    except Exception as exc:  # noqa: BLE001
        print(f"  [md-cleanup] post-processing failed: {exc!r}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Sync PDF and presentation files into Cursor-readable text files.")
    parser.add_argument(
        "--project-root",
        default=Path(__file__).resolve().parents[1],
        type=Path,
        help="Project root that contains the source documents.",
    )
    parser.add_argument(
        "--no-cleanup",
        action="store_true",
        help="Skip the md-cleanup post-processing step.",
    )
    args = parser.parse_args()

    project_root = args.project_root.resolve()
    manifests = run_sync(project_root)
    print(f"Synced {len(manifests)} document(s) into cursor_docs.")

    if not args.no_cleanup:
        config = load_config(project_root)
        output_dir = project_root / config["output_dir"]
        run_md_cleanup(output_dir)


if __name__ == "__main__":
    main()
